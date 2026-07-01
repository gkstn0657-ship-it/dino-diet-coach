# -*- coding: utf-8 -*-
import re, sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

prs = Presentation("C:/SSAFY/관통프로젝트/docs/냠냠코치_발표_초안_claude_v2.pptx")
EMOJI = re.compile("[\U0001F000-\U0001FAFF☀-➿]")
def clean(s): return EMOJI.sub("", s)

SW, SH = prs.slide_width, prs.slide_height
print("슬라이드 수:", len(prs.slides), " 캔버스:", round(SW/914400,2), "x", round(SH/914400,2), "in")
for i, sl in enumerate(prs.slides, 1):
    parts=[]
    overflow=[]
    for sh in sl.shapes:
        # out-of-bounds check
        try:
            if sh.left is not None and (sh.left < -9144 or sh.top < -9144 or
               sh.left+sh.width > SW+9144 or sh.top+sh.height > SH+9144):
                overflow.append(getattr(sh,'shape_type',''))
        except Exception:
            pass
        if sh.has_text_frame:
            t=clean(sh.text_frame.text).strip().replace("\n"," / ")
            if t: parts.append(t)
        if getattr(sh,'has_table',False):
            parts.append("[표]")
    line=" | ".join(parts)
    print(f"[{i}] {line[:180]}")
    if overflow:
        print(f"     !! 경계 이탈 의심 도형 {len(overflow)}개")
